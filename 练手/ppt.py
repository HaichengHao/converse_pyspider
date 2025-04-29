from pptx import Presentation
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import base64

# ===== 生成基础64卦数据 =====
hexagram_data = [
    {"name": "乾为天", "symbol": "䷀", "position": "世六应三", "usage": ["父母", "官鬼"]},
    {"name": "坤为地", "symbol": "䷁", "position": "世六应三", "usage": ["妻财", "子孙"]},
    {"name": "水雷屯", "symbol": "䷂", "position": "世四应初", "usage": ["兄弟", "官鬼"]},
    {"name": "山水蒙", "symbol": "䷃", "position": "世三应六", "usage": ["子孙", "父母"]},
    # 此处补充剩余60卦数据...
]

# ===== 创建演示文稿 =====
prs = Presentation()
prs.slide_width = Inches(13.33)  # 16:9宽屏
prs.slide_height = Inches(7.5)


# ===== 自定义版式 =====
def add_title_slide(text):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    title.text = text
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(59, 89, 152)


def add_content_slide(title_text, content):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title_text
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    for item in content:
        p = tf.add_paragraph()
        p.text = item
        p.level = 0


# ===== 封面页 =====
add_title_slide("六十四卦精要\n30分钟速通课程")

# ===== 核心概念页 =====
core_concepts = [
    "卦象结构：六爻组成，分上下卦",
    "世应定位：世为主，应为客，相差三位",
    "用神体系：六亲（父母/官鬼/兄弟/妻财/子孙）",
    "记忆口诀：一二三下卦，四五六上求"
]
add_content_slide("核心概念", core_concepts)


# ===== 卦象生成函数 =====
def create_hexagram_shape(slide, symbol, left, top):
    # 使用Unicode字符生成卦象
    symbol_box = slide.shapes.add_textbox(left, top, Inches(2), Inches(1))
    tf = symbol_box.text_frame
    p = tf.add_paragraph()
    p.text = symbol
    p.font.size = Pt(48)
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Segoe UI Symbol"


# ===== 重点卦象详解 =====
for idx, hexagram in enumerate(hexagram_data[:4]):  # 演示前4卦
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # 空白版式

    # 卦象显示
    create_hexagram_shape(slide, hexagram["symbol"], Inches(1), Inches(2))

    # 文字说明
    textbox = slide.shapes.add_textbox(Inches(4), Inches(2), Inches(6), Inches(4))
    tf = textbox.text_frame
    tf.word_wrap = True

    title = tf.add_paragraph()
    title.text = hexagram["name"]
    title.font.size = Pt(28)
    title.font.bold = True

    content = tf.add_paragraph()
    content.text = f"世应位置：{hexagram['position']}\n常用用神：{', '.join(hexagram['usage'])}"
    content.font.size = Pt(20)

    # 添加装饰线
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(3), Inches(3.5),
        Inches(3), Inches(4.5))
    line.line.width = Pt(2)

# ===== 记忆技巧页 =====
mnemonics = [
    "世爻定位歌诀：",
    "天同二世天变五，",
    "地同四世地变初，",
    "人同游魂人变归，",
    "本宫六世三世异。"
]
add_content_slide("记忆技巧", mnemonics)

# ===== 样式统一设置 =====
for slide in prs.slides:
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                paragraph.font.name = "微软雅黑"
                if shape == slide.shapes.title:
                    paragraph.font.size = Pt(36)
                    paragraph.font.color.rgb = RGBColor(198, 89, 17)
                else:
                    paragraph.font.size = Pt(20)

# ===== 保存文件 =====
prs.save('完整版六十四卦课程.pptx')
print("PPT生成完成！包含以下页面：")
print([slide.shapes.title.text for slide in prs.slides if slide.shapes.title])